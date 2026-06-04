"""V2 parser tests: DOCX, XLSX, PPTX, TXT, RTF.

Creates a sample file for each format in-memory, uploads it through the API and
asserts that text/tables/metadata are extracted correctly.
"""
from __future__ import annotations

import io
import os
import tempfile

os.environ["DATABASE_URL"] = f"sqlite:///{tempfile.gettempdir()}/udip_test_v2.db"

from fastapi.testclient import TestClient  # noqa: E402

import main  # noqa: E402
from app.database.session import init_db  # noqa: E402

init_db()
client = TestClient(main.app)


def _upload(name: str, data: bytes, mime: str) -> dict:
    r = client.post("/api/v1/upload", files={"file": (name, data, mime)})
    assert r.status_code == 200, r.text
    return r.json()["document"]


def _detail(public_id: str) -> dict:
    return client.get(f"/api/v1/documents/{public_id}").json()


# ---- DOCX ----------------------------------------------------------
def _make_docx() -> bytes:
    import docx

    d = docx.Document()
    d.add_heading("Shartnoma", level=1)
    d.add_paragraph("Bu test hujjati. To'lov muddati 30 kun.")
    table = d.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "Mahsulot"
    table.rows[0].cells[1].text = "Narx"
    table.rows[1].cells[0].text = "Kitob"
    table.rows[1].cells[1].text = "50000"
    buf = io.BytesIO()
    d.save(buf)
    return buf.getvalue()


def test_docx_extracts_text_and_table():
    doc = _upload("shartnoma.docx", _make_docx(),
                  "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
    assert doc["file_type"] == "docx"
    detail = _detail(doc["public_id"])
    text = detail["pages"][0]["text"]
    assert "Shartnoma" in text and "To'lov muddati" in text
    # Table is captured in layout blocks.
    blocks = detail["pages"][0]["layout"]["blocks"]
    assert any(b["type"] == "table" for b in blocks)


# ---- XLSX ----------------------------------------------------------
def _make_xlsx() -> bytes:
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Hisobot"
    ws.append(["Oy", "Summa"])
    ws.append(["Yanvar", 1000])
    ws.append(["Fevral", 2000])
    wb.create_sheet("Ikkinchi")
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def test_xlsx_reads_sheets():
    doc = _upload("hisobot.xlsx", _make_xlsx(),
                  "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
    detail = _detail(doc["public_id"])
    assert detail["page_count"] == 2  # two sheets
    assert "Yanvar" in detail["pages"][0]["text"]


# ---- PPTX ----------------------------------------------------------
def _make_pptx() -> bytes:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Taqdimot"
    slide.placeholders[1].text = "Birinchi slayd matni"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_pptx_reads_slides():
    doc = _upload("taqdimot.pptx", _make_pptx(),
                  "application/vnd.openxmlformats-officedocument.presentationml.presentation")
    detail = _detail(doc["public_id"])
    assert detail["page_count"] >= 1
    assert "Taqdimot" in detail["pages"][0]["text"]


# ---- TXT -----------------------------------------------------------
def test_txt_plain():
    doc = _upload("matn.txt", "Salom dunyo\nIkkinchi qator".encode("utf-8"), "text/plain")
    detail = _detail(doc["public_id"])
    assert "Salom dunyo" in detail["pages"][0]["text"]


# ---- RTF -----------------------------------------------------------
def test_rtf_stripped():
    rtf = r"{\rtf1\ansi\deff0 {\fonttbl {\f0 Times;}}\f0\fs24 Salom RTF hujjat.}"
    doc = _upload("hujjat.rtf", rtf.encode("utf-8"), "application/rtf")
    detail = _detail(doc["public_id"])
    assert "Salom RTF" in detail["pages"][0]["text"]


# ---- parse endpoint ------------------------------------------------
def test_parse_endpoint_returns_tables():
    doc = _upload("p.docx", _make_docx(),
                  "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
    r = client.post("/api/v1/parse", json={"document_id": doc["id"]})
    assert r.status_code == 200, r.text
    body = r.json()
    assert body["parser"] == "docx"
    assert len(body["tables"]) >= 1
    assert body["tables"][0]["rows"][0] == ["Mahsulot", "Narx"]
