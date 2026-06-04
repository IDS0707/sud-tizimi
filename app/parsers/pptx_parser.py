"""PowerPoint parser (TZ section 2.5).

Each slide becomes one ``ParsedPage``: shape text is collected as the slide
body, the title shape is flagged as a heading block, and speaker notes are
appended. Tables inside slides are captured as ``ParsedTable`` too.
"""
from __future__ import annotations

from app.parsers.base import BaseParser, ParsedPage, ParseResult, ParsedTable
from app.utils.logger import get_logger

log = get_logger("udip.parsers.pptx")

try:
    from pptx import Presentation

    _HAS_PPTX = True
except Exception:  # pragma: no cover
    _HAS_PPTX = False


class PptxParser(BaseParser):
    name = "pptx"
    extensions = ("pptx",)

    def parse(self, file_path: str, **kwargs) -> ParseResult:
        if not _HAS_PPTX:
            return ParseResult(parser=self.name, metadata={"error": "python-pptx not installed"})

        prs = Presentation(file_path)
        pages: list[ParsedPage] = []
        tables: list[ParsedTable] = []

        for idx, slide in enumerate(prs.slides, start=1):
            lines: list[str] = []
            blocks: list[dict] = []
            title = None

            for shape in slide.shapes:
                if shape.has_table:
                    rows = [[c.text.strip() for c in row.cells] for row in shape.table.rows]
                    if rows:
                        tables.append(ParsedTable(page_number=idx, rows=rows))
                        blocks.append({"type": "table", "rows": rows})
                        lines.extend("\t".join(r) for r in rows)
                    continue
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                is_title = (shape == slide.shapes.title)
                if is_title and title is None:
                    title = text
                    blocks.insert(0, {"type": "heading", "text": text})
                    lines.insert(0, "# " + text)
                else:
                    blocks.append({"type": "text", "text": text})
                    lines.append(text)

            # Speaker notes.
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    blocks.append({"type": "notes", "text": notes})
                    lines.append("\n[Izoh] " + notes)

            pages.append(ParsedPage(
                page_number=idx,
                text="\n".join(lines),
                blocks=blocks,
            ))

        meta = {"page_count": len(pages), "slides": len(pages), "tables": len(tables)}
        log.info("Parsed PPTX %s: %d slides", file_path, len(pages))
        return ParseResult(parser=self.name, pages=pages, tables=tables, metadata=meta)


pptx_parser = PptxParser()
